"""PPTX 导出器 - 将 markdown 内容转换为 PPTX"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def parse_markdown_slides(content: str) -> list[dict]:
    """解析 markdown 内容，提取各 slide"""
    slides = []
    current_slide = {"title": "", "content": []}

    lines = content.strip().split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 检测新 slide (以 # 开头)
        if line.startswith("#"):
            if current_slide["title"] or current_slide["content"]:
                slides.append(current_slide)
            current_slide = {"title": line.lstrip("#").strip(), "content": []}
        elif line.startswith("##"):
            current_slide["title"] = line.lstrip("#").strip()
        elif line.startswith("-"):
            current_slide["content"].append(
                {"type": "bullet", "text": line.lstrip("-").strip()}
            )
        elif line.startswith("**") and line.endswith("**"):
            current_slide["content"].append(
                {"type": "bold", "text": line.strip("**").strip()}
            )
        else:
            current_slide["content"].append({"type": "text", "text": line})

    # 添加最后一个 slide
    if current_slide["title"] or current_slide["content"]:
        slides.append(current_slide)

    return slides


def create_pptx(
    title: str,
    author: str,
    slides_data: list[dict],
    output_path: str,
    width_inches: float = 13.33,  # 16:9 宽屏
    height_inches: float = 7.5,
):
    """创建 PPTX 文件"""
    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)

    for i, slide_data in enumerate(slides_data):
        # 选择布局
        if i == 0:
            # 标题页
            slide_layout = prs.slide_layouts[6]  # 空白页
            slide = prs.slides.add_slide(slide_layout)

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(width_inches - 2), Inches(1.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", title)
            p.font.size = Pt(44)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # 副标题/作者
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(width_inches - 2), Inches(2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for item in slide_data.get("content", []):
                if item["type"] in ("text", "bold"):
                    p = tf.add_paragraph()
                    p.text = item["text"]
                    p.font.size = Pt(24)
                    if item["type"] == "bold":
                        p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

        else:
            # 内容页 - 标题 + 左侧文字 + 右侧占位
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(width_inches - 1), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Slide {i}: {slide_data.get('title', '')}"
            p.font.size = Pt(24)
            p.font.bold = True

            # 左侧内容区 (占左半边)
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(6), Inches(5.5)
            )
            tf = left_box.text_frame
            tf.word_wrap = True

            for item in slide_data.get("content", []):
                if item["type"] == "bullet":
                    p = tf.add_paragraph()
                    p.text = "• " + item["text"]
                    p.font.size = Pt(18)
                    p.space_before = Pt(12)
                elif item["type"] == "bold":
                    p = tf.add_paragraph()
                    p.text = item["text"]
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.space_before = Pt(18)
                else:
                    p = tf.add_paragraph()
                    p.text = item["text"]
                    p.font.size = Pt(16)

            # 右侧占位区 (占右半边)
            right_box = slide.shapes.add_textbox(
                Inches(7), Inches(1.3), Inches(5.8), Inches(5.5)
            )
            tf = right_box.text_frame
            tf.word_wrap = True

            p = tf.add_paragraph()
            p.text = "📊 图表占位 / Chart Placeholder"
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = RGBColor(128, 128, 128)
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = "\n描述 / Description:\n[在此处插入相关图表]\n\n建议:\n- 放入论文中最核心的 Figure\n- 保持简洁，突出重点"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(128, 128, 128)
            p.space_before = Pt(20)

    prs.save(output_path)
    return output_path


def export_to_pptx(
    markdown_content: str,
    paper_id: str,
    title: str,
    authors: str,
    output_dir: str = "./slides",
) -> str:
    """导出为 PPTX (便捷函数)"""
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    slides_data = parse_markdown_slides(markdown_content)

    # 第一个 slide 用标题，其余去掉前缀
    for i, slide in enumerate(slides_data):
        title_text = slide.get("title", "")
        # 移除 "Slide N: " 前缀用于显示
        if title_text.startswith("Slide "):
            parts = title_text.split(":", 1)
            if len(parts) > 1:
                slide["display_title"] = parts[1].strip()
            else:
                slide["display_title"] = title_text
        else:
            slide["display_title"] = title_text

    output_path = Path(output_dir) / f"{paper_id}_slide.pptx"
    create_pptx(
        title=title,
        author=authors,
        slides_data=slides_data,
        output_path=str(output_path),
    )

    return str(output_path)
